
import io
import mimetypes
from typing import Optional

# Kütüphaneleri ayrı ayrı import et, böylece biri eksikse diğeri çalışmaya devam eder
try:
    import fitz  # PyMuPDF
except ImportError:
    print("UYARI: 'pymupdf' kütüphanesi yüklü değil. PDF okuma çalışmayacak.")
    fitz = None

try:
    from docx import Document as DocxDocument
except ImportError:
    print("UYARI: 'python-docx' kütüphanesi yüklü değil. .docx okuma çalışmayacak.")
    DocxDocument = None

try:
    from openpyxl import load_workbook
except ImportError:
    print("UYARI: 'openpyxl' kütüphanesi yüklü değil. .xlsx okuma çalışmayacak.")
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    print("UYARI: 'python-pptx' kütüphanesi yüklü değil. .pptx okuma çalışmayacak.")
    Presentation = None
    
try:
    from google.cloud import vision
except ImportError:
    print("UYARI: 'google-cloud-vision' kütüphanesi yüklü değil. OCR çalışmayacak.")
    vision = None


def _read_text(file_bytes: bytes) -> str:
    """Düz metin dosyalarını (txt, csv, json vb.) okur."""
    try:
        return file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        try:
            return file_bytes.decode('latin-1') # Başka bir encoding dene
        except Exception as e:
            return f"[Metin dosyası okunurken hata: {str(e)}]"

def _read_pdf(file_bytes: bytes) -> str:
    """PDF dosyasının içeriğini okur. Görüntü tabanlı sayfalar için OCR dener."""
    if not fitz: return "[PDF okuyucu (PyMuPDF) yüklü değil. Lütfen 'pip install pymupdf' komutunu çalıştırın.]"
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text = ""
        page_count = len(doc)
        print(f"📄 PDF açıldı: {page_count} sayfa")
        
        for page_num, page in enumerate(doc):
            page_text = page.get_text()
            if page_text and page_text.strip():
                text += f"\n--- Sayfa {page_num + 1} ---\n{page_text}\n"
            else:
                # Eğer metin yoksa, görüntü olabilir - OCR denemesi yap
                print(f"⚠️ Sayfa {page_num + 1} metin içermiyor, görüntü olarak işleniyor...")
                try:
                    # Sayfayı görüntüye çevir ve OCR yap
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                    img_bytes = pix.tobytes("png")
                    ocr_text = _read_image_ocr_google_vision(img_bytes)
                    if ocr_text and not ocr_text.strip().startswith("["):
                        text += f"\n--- Sayfa {page_num + 1} (OCR) ---\n{ocr_text}\n"
                        print(f"✅ Sayfa {page_num + 1} OCR ile okundu: {len(ocr_text)} karakter")
                    else:
                        print(f"⚠️ Sayfa {page_num + 1} OCR ile de okunamadı")
                except Exception as ocr_error:
                    print(f"⚠️ Sayfa {page_num + 1} OCR hatası: {str(ocr_error)}")
        
        doc.close()
        print(f"📄 PDF okundu: {len(text)} karakter çıkarıldı")
        return text if text.strip() else "[PDF dosyası metin içermiyor veya görüntü tabanlı PDF. OCR gerekebilir.]"
    except Exception as e:
        print(f"❌ PDF okuma hatası: {str(e)}")
        return f"[PDF dosyası okunurken hata: {str(e)}]"

def _read_docx(file_bytes: bytes) -> str:
    """Word (.docx) dosyasının içeriğini okur."""
    if not DocxDocument: return "[Word okuyucu (python-docx) yüklü değil. Lütfen 'pip install python-docx' komutunu çalıştırın.]"
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        full_text = []

        # Önce paragraflardaki metinleri ekle
        for para in doc.paragraphs:
            if para.text:
                full_text.append(para.text)

        # Sonra tablolardaki metinleri ekle
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        full_text.append(cell.text)
        
        return "\n".join(full_text)
    except Exception as e:
        return f"[Word dosyası okunurken hata: {str(e)}]"

def _read_xlsx(file_bytes: bytes) -> str:
    """Excel (.xlsx) dosyasının içeriğini okur."""
    if not load_workbook: return "[Excel okuyucu (openpyxl) yüklü değil. Lütfen 'pip install openpyxl' komutunu çalıştırın.]"
    try:
        wb = load_workbook(filename=io.BytesIO(file_bytes), read_only=True)
        text = ""
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"--- Sayfa: {sheet_name} ---\n"
            for row in sheet.iter_rows():
                row_text = "\t".join([str(cell.value) for cell in row if cell.value is not None])
                if row_text:
                    text += row_text + "\n"
        return text
    except Exception as e:
        return f"[Excel dosyası okunurken hata: {str(e)}]"

def _read_pptx(file_bytes: bytes) -> str:
    """PowerPoint (.pptx) dosyasının içeriğini okur."""
    if not Presentation: return "[PowerPoint okuyucu (python-pptx) yüklü değil. Lütfen 'pip install python-pptx' komutunu çalıştırın.]"
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        slide_text += para.text + "\n"
            if slide_text:
                 text += f"--- Slayt {prs.slides.index(slide) + 1} ---\n{slide_text}\n"
        return text
    except Exception as e:
        return f"[PowerPoint dosyası okunurken hata: {str(e)}]"

def _read_image_ocr_google_vision(file_bytes: bytes) -> str:
    """Resim dosyalarından Google Cloud Vision API ile metin okur."""
    if not vision:
        print("⚠️ Google Cloud Vision kütüphanesi yüklü değil")
        return "[Google Cloud Vision kütüphanesi yüklü değil. Lütfen 'pip install google-cloud-vision' komutunu çalıştırın.]"
    
    try:
        print(f"🖼️ Resim OCR başlatılıyor: {len(file_bytes)} bytes")
        client = vision.ImageAnnotatorClient()
        image = vision.Image(content=file_bytes)
        response = client.document_text_detection(image=image)
        
        if response.error.message:
            raise Exception(f"Cloud Vision API Hatası: {response.error.message}")
            
        if response.full_text_annotation:
            text = response.full_text_annotation.text
            print(f"✅ OCR tamamlandı: {len(text)} karakter çıkarıldı")
            return text
        else:
            print("⚠️ OCR sonucu boş - resimde metin bulunamadı")
            return "[Resimden metin (Cloud Vision) okunamadı]"
    
    except Exception as e:
        print(f"❌ OCR hatası: {str(e)}")
        return f"[Resim (Cloud Vision OCR) okunurken hata: {str(e)}]"


def extract_text_from_file(
    file_bytes: bytes, 
    file_name: str,
    mime_type: Optional[str] = None
) -> str:
    """
    Ana yönlendirici fonksiyon.
    (MIME tipine göre yönlendirme, dosya uzantısına göre fallback)
    """
    
    # Dosya uzantısını al (küçük harfe çevir)
    file_ext = file_name.lower().split('.')[-1] if '.' in file_name else ''
    
    # 1. MIME type hiç yoksa veya generic ise (octet-stream), dosya adından tahmin et
    if not mime_type or mime_type == "application/octet-stream":
        guessed_type, _ = mimetypes.guess_type(file_name)
        if guessed_type:
            print(f"⚠️ MIME type '{mime_type}' yetersiz, dosya adından tahmin edildi: {guessed_type}")
            mime_type = guessed_type

    print(f"Dosya okunuyor: {file_name} (MIME: {mime_type}, Uzantı: .{file_ext})")

    # ÖNCE: Dosya uzantısına göre kontrol et (MIME type yanlış olabilir)
    if file_ext == 'docx':
        print(f"📄 DOCX uzantısı tespit edildi, DOCX parser'a yönlendiriliyor...")
        result = _read_docx(file_bytes)
        # Eğer parser başarılıysa (hata mesajı yoksa), sonucu döndür
        if result and not result.startswith("[Word"):
            return result
        # Parser başarısız olduysa (kütüphane yoksa veya hata varsa), MIME type kontrolüne geç
        print(f"⚠️ DOCX parser sonucu: {result[:100]}...")
    
    if file_ext == 'xlsx':
        print(f"📊 XLSX uzantısı tespit edildi, XLSX parser'a yönlendiriliyor...")
        result = _read_xlsx(file_bytes)
        if result and not result.startswith("[Excel"):
            return result
    
    if file_ext == 'pptx':
        print(f"📽️ PPTX uzantısı tespit edildi, PPTX parser'a yönlendiriliyor...")
        result = _read_pptx(file_bytes)
        if result and not result.startswith("[PowerPoint"):
            return result
    
    if file_ext == 'pdf':
        print(f"📄 PDF uzantısı tespit edildi, PDF parser'a yönlendiriliyor...")
        return _read_pdf(file_bytes)

    # MIME tipine göre yönlendir
    if not mime_type:
        return _read_text(file_bytes)
    
    if mime_type.startswith("text/"):
        return _read_text(file_bytes)
        
    elif mime_type == "application/pdf":
        return _read_pdf(file_bytes)
        
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _read_docx(file_bytes)
        
    elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return _read_xlsx(file_bytes)
        
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _read_pptx(file_bytes)
        
    elif mime_type.startswith("image/"):
        return _read_image_ocr_google_vision(file_bytes)

    # Eğer yukarıdakilerden hiçbiri değilse, son bir şans dosya uzantısına tekrar bak
    # Bazı tarayıcılar docx için farklı mime tipleri gönderebilir
    else:
        guessed_type, _ = mimetypes.guess_type(file_name)
        if guessed_type and guessed_type != mime_type:
            print(f"⚠️ Desteklenmeyen MIME '{mime_type}', dosya uzantısına göre tekrar deneniyor: {guessed_type}")
            # Recursive call with the guessed type
            return extract_text_from_file(file_bytes, file_name, guessed_type)
        
        # Son çare: Dosya uzantısına göre direkt kontrol
        if file_ext in ['docx', 'doc']:
            print(f"⚠️ MIME type desteklenmiyor ama .{file_ext} uzantısı var, DOCX parser deneniyor...")
            return _read_docx(file_bytes)
        elif file_ext in ['xlsx', 'xls']:
            print(f"⚠️ MIME type desteklenmiyor ama .{file_ext} uzantısı var, XLSX parser deneniyor...")
            return _read_xlsx(file_bytes)
        elif file_ext in ['pptx', 'ppt']:
            print(f"⚠️ MIME type desteklenmiyor ama .{file_ext} uzantısı var, PPTX parser deneniyor...")
            return _read_pptx(file_bytes)
            
        return f"[Desteklenmeyen dosya formatı: {mime_type} (uzantı: .{file_ext})]"