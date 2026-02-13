# backend/app/services/presentation_agent_service.py
# Sunum oluşturma için özel agent servisi

import io
import json
import tempfile
import os
from typing import Dict, Any, Optional, List, Tuple
from app.services.llm_providers import get_llm_for_model

# python-pptx import'u
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    print("UYARI: 'python-pptx' kütüphanesi yüklü değil. Sunum oluşturma çalışmayacak.")
    PPTX_AVAILABLE = False
    Presentation = None

ANALYSIS_PROMPT = """Sen bir sunum hazırlama uzmanısın. Kullanıcı bir konu hakkında sunum hazırlamak istiyor.

KULLANICI KONUSU: {topic}

MEVCUT BİLGİLER:
{context_info}

{user_answered_note}

GÖREVİN:
1. Mevcut bilgileri analiz et
2. {analysis_instruction}
3. Eğer eksik bilgiler varsa, bunları kullanıcıya sormak için JSON formatında döndür
4. Eğer yeterli bilgi varsa, sunum yapısını oluştur

ÇIKTI FORMATI (JSON):
Eğer eksik bilgiler varsa:
{{
    "status": "missing_info",
    "missing_fields": [
        {{
            "field": "Hedef kitle",
            "question": "Bu sunumun hedef kitlesi kimler? (ör: yönetim kurulu, müşteriler, çalışanlar)",
            "required": true
        }},
        {{
            "field": "Sunum süresi",
            "question": "Sunum ne kadar sürecek? (ör: 15 dakika, 30 dakika)",
            "required": false
        }}
    ]
}}

Eğer yeterli bilgi varsa:
{{
    "status": "ready",
    "presentation_structure": {{
        "title": "Sunum Başlığı",
        "subtitle": "Alt başlık (opsiyonel)",
        "slides": [
            {{
                "slide_number": 1,
                "slide_type": "title",
                "title": "Başlık",
                "content": []
            }},
            {{
                "slide_number": 2,
                "slide_type": "content",
                "title": "Slayt Başlığı",
                "content": [
                    "Ana nokta 1",
                    "Ana nokta 2",
                    "Ana nokta 3"
                ]
            }}
        ]
    }}
}}

Lütfen analiz yap ve JSON formatında cevap ver. Sadece JSON döndür, başka açıklama ekleme.
"""

PRESENTATION_PROMPT = """Sen bir sunum içerik uzmanısın. Verilen yapıya göre detaylı sunum içeriği oluştur.

SUNUM YAPISI:
{structure}

GÖREVİN:
Her slayt için detaylı içerik oluştur. Her slayt için:
- Başlık (structure'dan al)
- Ana noktalar (3-5 madde)
- Her madde için kısa açıklama (1-2 cümle)

ÖNEMLİ: JSON yanıtını tamamla. Kesilmiş JSON gönderme. Tüm slaytları içeren tam bir JSON döndür.

ÇIKTI FORMATI (JSON):
{{
    "slides": [
        {{
            "slide_number": 1,
            "title": "Slayt Başlığı",
            "bullet_points": [
                {{
                    "point": "Ana nokta",
                    "description": "Kısa açıklama"
                }}
            ]
        }}
    ]
}}

Lütfen sadece JSON döndür, başka açıklama ekleme. JSON'u tamamla, kesme.
"""


def analyze_presentation_requirements(
    topic: str,
    context_info: str = "",
    model_name: str = "gemini",
    user_answered_questions: bool = False
) -> Dict[str, Any]:
    """Sunum gereksinimlerini analiz eder ve eksik bilgileri tespit eder."""
    try:
        llm = get_llm_for_model(model_name)
        
        # Kullanıcı soruları cevapladıysa, daha net bir talimat ver
        if user_answered_questions:
            user_answered_note = "ÖNEMLİ: Kullanıcı daha önce sorduğun soruları cevapladı. Artık yeterli bilgiye sahipsin, sunum yapısını oluşturmalısın."
            analysis_instruction = "Kullanıcının verdiği cevapları dikkate alarak sunum için yeterli bilgi olup olmadığını kontrol et"
        else:
            user_answered_note = ""
            analysis_instruction = "Sunum için gerekli olan ama eksik olan bilgileri tespit et"
        
        prompt = ANALYSIS_PROMPT.format(
            topic=topic,
            context_info=context_info if context_info else "Kullanıcı henüz bağlam dosyası eklememiş.",
            user_answered_note=user_answered_note,
            analysis_instruction=analysis_instruction
        )
        
        response = llm.invoke(prompt)
        response_text = response.content if hasattr(response, 'content') else str(response)
        
        # JSON'u temizle (markdown code block varsa kaldır)
        response_text = response_text.strip()
        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.startswith("```"):
            response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        response_text = response_text.strip()
        
        result = json.loads(response_text)
        return result
        
    except json.JSONDecodeError as e:
        print(f"❌ JSON parse hatası: {e}")
        print(f"LLM yanıtı: {response_text[:500]}")
        return {
            "status": "error",
            "message": "Sunum analizi sırasında bir hata oluştu. Lütfen tekrar deneyin."
        }
    except Exception as e:
        print(f"❌ Sunum analizi hatası: {e}")
        import traceback
        traceback.print_exc()
        return {
            "status": "error",
            "message": f"Sunum analizi sırasında bir hata oluştu: {str(e)}"
        }


def generate_presentation_content(
    structure: Dict[str, Any],
    model_name: str = "gemini"
) -> Dict[str, Any]:
    """Sunum yapısına göre detaylı içerik oluşturur."""
    try:
        llm = get_llm_for_model(model_name)
        prompt = PRESENTATION_PROMPT.format(
            structure=json.dumps(structure, ensure_ascii=False, indent=2)
        )
        
        response = llm.invoke(prompt)
        response_text = response.content if hasattr(response, 'content') else str(response)
        
        # JSON'u temizle
        response_text = response_text.strip()
        if response_text.startswith("```json"):
            response_text = response_text[7:]
        if response_text.startswith("```"):
            response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        response_text = response_text.strip()
        
        # JSON parse dene
        try:
            result = json.loads(response_text)
            return result
        except json.JSONDecodeError:
            # JSON kesilmiş olabilir, düzeltmeye çalış
            print(f"⚠️ JSON parse hatası, düzeltmeye çalışılıyor...")
            # Kapanmamış string'leri kapat
            if response_text.count('"') % 2 != 0:
                # Tek sayıda tırnak var, son tırnağı kapat
                last_quote = response_text.rfind('"')
                if last_quote > 0 and response_text[last_quote-1] != '\\':
                    response_text = response_text[:last_quote+1] + '"'
            
            # Kapanmamış objeleri kapat
            open_braces = response_text.count('{')
            close_braces = response_text.count('}')
            if open_braces > close_braces:
                response_text += '}' * (open_braces - close_braces)
            
            open_brackets = response_text.count('[')
            close_brackets = response_text.count(']')
            if open_brackets > close_brackets:
                response_text += ']' * (open_brackets - close_brackets)
            
            try:
                result = json.loads(response_text)
                print(f"✅ JSON düzeltildi ve parse edildi")
                return result
            except json.JSONDecodeError as e2:
                print(f"❌ JSON parse hatası (düzeltme sonrası): {e2}")
                print(f"LLM yanıtı (ilk 2000 karakter): {response_text[:2000]}")
                
                # Daha agresif düzeltme: Kesilmiş JSON'u bul ve parse et
                # Son geçerli JSON objesini bul
                try:
                    # Son kapanan objeyi bul
                    last_valid_brace = response_text.rfind('}')
                    if last_valid_brace > 0:
                        # Son geçerli objeye kadar al
                        partial_json = response_text[:last_valid_brace + 1]
                        # Kapanmamış array'leri kapat
                        open_brackets = partial_json.count('[')
                        close_brackets = partial_json.count(']')
                        if open_brackets > close_brackets:
                            partial_json += ']' * (open_brackets - close_brackets)
                        # Kapanmamış objeleri kapat
                        open_braces = partial_json.count('{')
                        close_braces = partial_json.count('}')
                        if open_braces > close_braces:
                            partial_json += '}' * (open_braces - close_braces)
                        
                        # Son kapanan array'i bul
                        last_valid_bracket = partial_json.rfind(']')
                        if last_valid_bracket > 0:
                            partial_json = partial_json[:last_valid_bracket + 1]
                            # Kapanmamış objeleri kapat
                            open_braces = partial_json.count('{')
                            close_braces = partial_json.count('}')
                            if open_braces > close_braces:
                                partial_json += '}' * (open_braces - close_braces)
                        
                        # JSON'u tamamla
                        if not partial_json.strip().endswith('}'):
                            partial_json = '{"slides": ' + partial_json + '}'
                        else:
                            # Zaten tamamlanmış görünüyor
                            pass
                        
                        result = json.loads(partial_json)
                        print(f"✅ Kesilmiş JSON'dan kısmi veri çıkarıldı: {len(result.get('slides', []))} slayt")
                        return result
                except:
                    pass
                
                # En azından boş bir yapı döndür
                return {"slides": []}
        
    except json.JSONDecodeError as e:
        print(f"❌ JSON parse hatası: {e}")
        print(f"LLM yanıtı (ilk 1000 karakter): {response_text[:1000] if 'response_text' in locals() else 'Yanıt alınamadı'}")
        return {"slides": []}
    except Exception as e:
        print(f"❌ İçerik oluşturma hatası: {e}")
        import traceback
        traceback.print_exc()
        return {"slides": []}


def create_presentation_file(
    presentation_data: Dict[str, Any],
    title: str = "Sunum"
) -> bytes:
    """PowerPoint sunumu oluşturur ve bytes olarak döndürür."""
    if not PPTX_AVAILABLE:
        raise Exception("python-pptx kütüphanesi yüklü değil.")
    
    try:
        # Yeni sunum oluştur
        prs = Presentation()
        
        # Slayt genişliği ve yüksekliği (16:9)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        slides_data = presentation_data.get("slides", [])
        
        # Eğer slides boşsa, en azından başlık slaytı oluştur
        if not slides_data:
            print(f"⚠️ Slayt verisi boş, sadece başlık slaytı oluşturuluyor")
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = presentation_data.get("title", "Sunum")
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)
            title_shape.text_frame.paragraphs[0].font.bold = True
            # Sunumu kaydet ve döndür
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            return output.getvalue()
        
        for slide_info in slides_data:
            slide_number = slide_info.get("slide_number", 0)
            slide_type = slide_info.get("slide_type", "content")
            slide_title = slide_info.get("title", "")
            bullet_points = slide_info.get("bullet_points", [])
            
            # Slayt oluştur
            if slide_type == "title":
                # Başlık slaytı
                slide_layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(slide_layout)
                
                title_shape = slide.shapes.title
                title_shape.text = slide_title
                title_shape.text_frame.paragraphs[0].font.size = Pt(44)
                title_shape.text_frame.paragraphs[0].font.bold = True
                
                # Alt başlık varsa
                if len(slide.shapes.placeholders) > 1:
                    subtitle = slide.shapes.placeholders[1]
                    subtitle_text = presentation_data.get("subtitle", "")
                    if subtitle_text:
                        subtitle.text = subtitle_text
                        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            else:
                # İçerik slaytı
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Başlık
                title_shape = slide.shapes.title
                title_shape.text = slide_title
                title_shape.text_frame.paragraphs[0].font.size = Pt(32)
                title_shape.text_frame.paragraphs[0].font.bold = True
                
                # İçerik
                if len(slide.shapes.placeholders) > 1:
                    content_shape = slide.shapes.placeholders[1]
                    text_frame = content_shape.text_frame
                    text_frame.word_wrap = True
                    
                    # Mevcut paragrafı temizle
                    text_frame.clear()
                    
                    # Her bullet point için paragraf ekle
                    for idx, point_data in enumerate(bullet_points):
                        point = point_data.get("point", "")
                        description = point_data.get("description", "")
                        
                        if idx == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = point
                        p.font.size = Pt(18)
                        p.font.bold = True
                        p.level = 0
                        
                        # Açıklama varsa alt paragraf olarak ekle
                        if description:
                            sub_p = text_frame.add_paragraph()
                            sub_p.text = description
                            sub_p.font.size = Pt(14)
                            sub_p.level = 1
        
        # Sunumu bytes'a çevir
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()
        
    except Exception as e:
        print(f"❌ Sunum oluşturma hatası: {e}")
        import traceback
        traceback.print_exc()
        raise Exception(f"Sunum oluşturulurken hata oluştu: {str(e)}")


def extract_context_info(context_files: List[Dict[str, Any]], db, storage, user) -> str:
    """Context dosyalarından bilgi çıkarır."""
    if not context_files:
        return ""
    
    context_texts = []
    from app.core.parsers import extract_text_from_file
    
    for context_file in context_files:
        if context_file.get("type") == "file":
            try:
                file_record = db.get_file_by_id(user.tenant_id, context_file["id"])
                if file_record:
                    # Eğer dosya external storage'dan geliyorsa, Google Drive/OneDrive'dan indir
                    if file_record.external_file_id and file_record.external_storage_type:
                        from google.cloud import firestore
                        from app.storage_adapters.google_drive_adapter import GoogleDriveAdapter
                        from app.storage_adapters.onedrive_adapter import OneDriveAdapter
                        from app.core.config import GOOGLE_DRIVE_CLIENT_ID, GOOGLE_DRIVE_CLIENT_SECRET, ONEDRIVE_CLIENT_ID, ONEDRIVE_CLIENT_SECRET
                        
                        firestore_db = firestore.Client()
                        storage_type = file_record.external_storage_type
                        
                        # Kullanıcının storage bağlantısını al
                        if storage_type == "google_drive":
                            user_storage = firestore_db.collection("user_external_storage").document(user.id).get()
                            if not user_storage.exists:
                                admin_settings = firestore_db.collection("external_storage_settings").document(user.tenant_id).get()
                                if not admin_settings.exists:
                                    continue
                                admin_data = admin_settings.to_dict()
                                access_token = admin_data.get('google_drive_access_token')
                                refresh_token = admin_data.get('google_drive_refresh_token')
                                client_id = GOOGLE_DRIVE_CLIENT_ID
                                client_secret = GOOGLE_DRIVE_CLIENT_SECRET
                            else:
                                storage_data = user_storage.to_dict()
                                access_token = storage_data.get('access_token')
                                refresh_token = storage_data.get('refresh_token')
                                client_id = GOOGLE_DRIVE_CLIENT_ID
                                client_secret = GOOGLE_DRIVE_CLIENT_SECRET
                            
                            adapter = GoogleDriveAdapter()
                        elif storage_type == "onedrive":
                            user_storage = firestore_db.collection("user_external_storage").document(user.id).get()
                            if not user_storage.exists:
                                admin_settings = firestore_db.collection("external_storage_settings").document(user.tenant_id).get()
                                if not admin_settings.exists:
                                    continue
                                admin_data = admin_settings.to_dict()
                                access_token = admin_data.get('onedrive_access_token')
                                refresh_token = admin_data.get('onedrive_refresh_token')
                                client_id = ONEDRIVE_CLIENT_ID
                                client_secret = ONEDRIVE_CLIENT_SECRET
                            else:
                                storage_data = user_storage.to_dict()
                                access_token = storage_data.get('access_token')
                                refresh_token = storage_data.get('refresh_token')
                                client_id = ONEDRIVE_CLIENT_ID
                                client_secret = ONEDRIVE_CLIENT_SECRET
                            
                            adapter = OneDriveAdapter()
                        else:
                            continue
                        
                        if not access_token:
                            continue
                        
                        # Token'ı kontrol et ve gerekirse yenile
                        try:
                            if storage_type == "google_drive":
                                file_bytes = adapter.download_file(
                                    file_id=file_record.external_file_id,
                                    access_token=access_token,
                                    mime_type=file_record.content_type
                                )
                            else:
                                file_bytes = adapter.download_file(
                                    file_id=file_record.external_file_id,
                                    access_token=access_token
                                )
                        except Exception as e:
                            if refresh_token and client_id and client_secret:
                                try:
                                    tokens = adapter.refresh_access_token(
                                        refresh_token=refresh_token,
                                        client_id=client_id,
                                        client_secret=client_secret
                                    )
                                    access_token = tokens['access_token']
                                    
                                    if user_storage.exists:
                                        firestore_db.collection("user_external_storage").document(user.id).update({
                                            'access_token': access_token
                                        })
                                    else:
                                        update_data = {}
                                        if storage_type == "google_drive":
                                            update_data['google_drive_access_token'] = access_token
                                        else:
                                            update_data['onedrive_access_token'] = access_token
                                        firestore_db.collection("external_storage_settings").document(user.tenant_id).update(update_data)
                                    
                                    if storage_type == "google_drive":
                                        file_bytes = adapter.download_file(
                                            file_id=file_record.external_file_id,
                                            access_token=access_token,
                                            mime_type=file_record.content_type
                                        )
                                    else:
                                        file_bytes = adapter.download_file(
                                            file_id=file_record.external_file_id,
                                            access_token=access_token
                                        )
                                except:
                                    continue
                            else:
                                continue
                    else:
                        # Normal dosyalar için mevcut mantık
                        if not file_record.storage_path:
                            continue
                        file_bytes = storage.download_file_content(file_record.storage_path)
                    
                    text = extract_text_from_file(
                        file_bytes=file_bytes,
                        file_name=file_record.name,
                        mime_type=file_record.content_type
                    )
                    context_texts.append(f"--- Dosya: {file_record.name} ---\n{text[:2000]}\n")  # İlk 2000 karakter
            except Exception as e:
                print(f"⚠️ Context dosyası okunamadı: {e}")
    
    return "\n".join(context_texts)

